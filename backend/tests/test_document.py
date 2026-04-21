import io

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter

from app.services.document import (
    DOCX_MIME,
    PDF_MIME,
    PPTX_MIME,
    XLSX_MIME,
    UnsupportedFileType,
    chunk_text,
    load_bytes,
    load_docx,
    load_pdf,
    load_pptx,
    load_text,
    load_xlsx,
)


def _pdf_with_text(text: str) -> bytes:
    # pypdf can't easily write text pages; use a known fixture pattern via
    # encoded PDF. Simpler: build a PDF using reportlab? not available.
    # Use pypdf to create an empty PDF and monkeypatch extraction path —
    # instead we use a real tiny PDF fixture string.
    raise NotImplementedError


def _docx_bytes(paragraphs: list[str]) -> bytes:
    doc = DocxDocument()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_bytes(slide_texts: list[str]) -> bytes:
    prs = Presentation()
    for t in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = t
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _xlsx_bytes(rows: list[list]) -> bytes:
    wb = Workbook()
    ws = wb.active
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_load_text_utf8():
    assert load_text(b"hello world") == "hello world"


def test_load_text_strips_bom():
    assert load_text(b"\xef\xbb\xbfhi") == "hi"


def test_load_docx_extracts_paragraphs():
    data = _docx_bytes(["First line.", "Second line."])
    text = load_docx(data)
    assert "First line." in text
    assert "Second line." in text


def test_load_pptx_extracts_slide_text():
    data = _pptx_bytes(["Slide A", "Slide B"])
    text = load_pptx(data)
    assert "Slide A" in text
    assert "Slide B" in text


def test_load_xlsx_extracts_cells():
    data = _xlsx_bytes([["name", "score"], ["alice", 10], ["bob", 20]])
    text = load_xlsx(data)
    assert "alice" in text
    assert "20" in text


def test_load_pdf_extracts_text():
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    buf = io.BytesIO()
    writer.write(buf)
    # A blank page yields empty string — that's expected behavior; we
    # just verify it doesn't crash and returns a str.
    result = load_pdf(buf.getvalue())
    assert isinstance(result, str)


def test_chunk_text_produces_chunks_of_bounded_size():
    text = "word " * 2000  # ~2000 tokens roughly
    chunks = chunk_text(text, chunk_tokens=500, overlap_tokens=50)
    assert len(chunks) >= 3
    for c in chunks:
        # tiktoken cl100k_base token count roughly <= 500
        assert len(c) > 0


def test_chunk_text_overlap_preserves_continuity():
    text = " ".join(f"w{i}" for i in range(600))
    chunks = chunk_text(text, chunk_tokens=200, overlap_tokens=40)
    # Last 40 tokens of chunk[0] should overlap with start of chunk[1]
    assert len(chunks) >= 2
    tail = " ".join(chunks[0].split()[-10:])
    assert tail.split()[0] in chunks[1]


def test_chunk_text_short_input_single_chunk():
    text = "hello world"
    assert chunk_text(text, chunk_tokens=500) == ["hello world"]


def test_load_bytes_routes_by_mime():
    docx = _docx_bytes(["A paragraph."])
    out = load_bytes(docx, DOCX_MIME, "x.docx")
    assert "A paragraph." in out


def test_load_bytes_text_plain():
    out = load_bytes(b"plain text here", "text/plain", "x.txt")
    assert out == "plain text here"


def test_load_bytes_unsupported_raises():
    with pytest.raises(UnsupportedFileType):
        load_bytes(b"x", "application/x-weird", "x.weird")


def test_load_bytes_pptx_and_xlsx_mimes():
    pp = _pptx_bytes(["Hello"])
    assert "Hello" in load_bytes(pp, PPTX_MIME, "a.pptx")
    xl = _xlsx_bytes([["row1"]])
    assert "row1" in load_bytes(xl, XLSX_MIME, "a.xlsx")

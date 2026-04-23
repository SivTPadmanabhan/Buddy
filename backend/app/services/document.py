import io
from dataclasses import dataclass
from typing import Callable

from docx import Document as DocxDocument
from langchain_experimental.text_splitter import SemanticChunker
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from app.logging_config import get_logger
from app.services.embeddings import Embedder
from app.services.langchain_embeddings import LangChainEmbeddingsAdapter
from app.services.text_preprocessor import preprocess_for_chunking

log = get_logger("app.services.document")

PDF_MIME = "application/pdf"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
TEXT_MIMES = {"text/plain", "text/csv", "text/markdown"}
IMAGE_MIMES = {"image/png", "image/jpeg", "image/jpg"}


class UnsupportedFileType(Exception):
    pass


@dataclass
class ChunkResult:
    text: str
    vector: list[float]
    chunk_index: int


def load_text(content: bytes) -> str:
    return content.decode("utf-8", errors="replace").lstrip("﻿")


def load_pdf(content: bytes) -> str:
    reader = PdfReader(io.BytesIO(content))
    parts = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts).strip()


def load_docx(content: bytes) -> str:
    doc = DocxDocument(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def load_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs) or para.text
                    if txt.strip():
                        parts.append(txt)
    return "\n".join(parts)


def load_xlsx(content: bytes) -> str:
    wb = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    lines = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append("\t".join(cells))
    return "\n".join(lines)


def semantic_chunk(
    text: str,
    mime_type: str,
    embedder: Embedder,
    breakpoint_type: str = "percentile",
    breakpoint_threshold: float = 90.0,
    structured_chunk_size: int = 500,
    structured_chunk_overlap: int = 50,
) -> list[ChunkResult]:
    text = text.strip()
    if not text:
        return []

    blocks = preprocess_for_chunking(text, mime_type)
    if not blocks:
        return []

    adapter = LangChainEmbeddingsAdapter(embedder)
    results: list[ChunkResult] = []
    index = 0

    for block in blocks:
        if block.skip_semantic:
            separators = ["\n"] if block.block_type == "table" else ["\n\n", "\n", " "]
            splitter = RecursiveCharacterTextSplitter(
                separators=separators,
                chunk_size=structured_chunk_size,
                chunk_overlap=structured_chunk_overlap,
                length_function=len,
            )
            chunk_texts = splitter.split_text(block.text)
        else:
            chunker = SemanticChunker(
                embeddings=adapter,
                breakpoint_threshold_type=breakpoint_type,
                breakpoint_threshold_amount=breakpoint_threshold,
            )
            docs = chunker.create_documents([block.text])
            chunk_texts = [doc.page_content for doc in docs]

        if not chunk_texts:
            continue

        vectors = embedder.embed_batch(chunk_texts)
        for ct, vector in zip(chunk_texts, vectors):
            results.append(ChunkResult(
                text=ct,
                vector=vector,
                chunk_index=index,
            ))
            index += 1

    return results


def _import_ocr() -> Callable[[bytes], str]:
    from app.services.ocr import extract_text
    return extract_text


_LOADERS: dict[str, Callable[[bytes], str]] = {
    PDF_MIME: load_pdf,
    DOCX_MIME: load_docx,
    PPTX_MIME: load_pptx,
    XLSX_MIME: load_xlsx,
}


def load_bytes(content: bytes, mime_type: str, filename: str) -> str:
    loader = _LOADERS.get(mime_type)
    if loader:
        return loader(content)
    if mime_type in TEXT_MIMES:
        return load_text(content)
    if mime_type in IMAGE_MIMES:
        return _import_ocr()(content)
    log.warning(
        "unsupported_file_type",
        category="sync",
        action="parse_document",
        mime_type=mime_type,
        filename=filename,
    )
    raise UnsupportedFileType(f"{mime_type} ({filename})")
